from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree
from pptx.oxml.ns import qn
def rgb_to_hex(rgb_color):
    return '{:02X}{:02X}{:02X}'.format(rgb_color[0], rgb_color[1], rgb_color[2])


def create_rounded_rectangle(shapes, attribute, **kwargs):

    left, top, width, height = attribute
    rounded_rect = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    rounded_rect.adjustments[0] = kwargs.get('adjustment', 0)

    colors = kwargs.get('colors')
    if colors and kwargs.get('gradient'):
        # Clean any existing fill
        fixed_stops = [
            (0, (112, 48, 160)),     # 0%
            (32000, (246, 146, 0)),  # 32%
            (60000, (254, 195, 0)),  # 60%
            (100000, (255, 255, 255))# 100%
        ]

        # Remove existing fill
        sp = rounded_rect._element
        fill = sp.find(qn('p:spPr'))
        for node in list(fill):
            if node.tag.endswith('solidFill') or node.tag.endswith('gradFill'):
                fill.remove(node)

        # Build the gradient XML
        stops_xml = ""
        for pos, rgb in fixed_stops:
            hex_color = rgb_to_hex(rgb)
            stops_xml += f'''
            <a:gs pos="{pos}">
                <a:srgbClr val="{hex_color}"/>
            </a:gs>
            '''

        gradient_xml = f'''
        <a:gradFill flip="none" rotWithShape="1" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:gsLst>
                {stops_xml}
            </a:gsLst>
            <a:lin ang="0" scaled="1"/>
        </a:gradFill>
        '''
        grad_fill = etree.fromstring(gradient_xml)
        fill.append(grad_fill)
    elif kwargs.get('color'):
        rounded_rect.fill.solid()
        rounded_rect.fill.fore_color.rgb = kwargs.get('color')

    # Line color
    rounded_rect.line.color.rgb = kwargs.get('line_color', RGBColor(255, 255, 255))

    # Texts
    if kwargs.get('texts'):
        text_frame = rounded_rect.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.alignment = kwargs.get('text_alignment', 1)
        run = p.add_run()
        run.text = kwargs.get('texts')
        run.font.size = kwargs.get('fontsize', Pt(6))
        run.font.color.rgb = kwargs.get('text_color', (0, 0, 0))

    # Shadow
    if not kwargs.get('shadow'):
        rounded_rect.shadow.inherit = False
        rounded_rect.shadow.visible = True

    return rounded_rect


def create_rectangle(shapes, attribute, **kwargs):

    left, top, width, height = attribute
    rectangle = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    colors = kwargs.get('colors')
    if colors and kwargs.get('gradient'):
        # Clean any existing fill
        fixed_stops = [
            (0, (112, 48, 160)),     # 0%
            (32000, (246, 146, 0)),  # 32%
            (60000, (254, 195, 0)),  # 60%
            (100000, (255, 255, 255))# 100%
        ]

        # Remove existing fill
        sp = rectangle._element
        fill = sp.find(qn('p:spPr'))
        for node in list(fill):
            if node.tag.endswith('solidFill') or node.tag.endswith('gradFill'):
                fill.remove(node)

        # Build the gradient XML
        stops_xml = ""
        for pos, rgb in fixed_stops:
            hex_color = rgb_to_hex(rgb)
            stops_xml += f'''
            <a:gs pos="{pos}">
                <a:srgbClr val="{hex_color}"/>
            </a:gs>
            '''

        gradient_xml = f'''
        <a:gradFill flip="none" rotWithShape="1" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:gsLst>
                {stops_xml}
            </a:gsLst>
            <a:lin ang="0" scaled="1"/>
        </a:gradFill>
        '''
        grad_fill = etree.fromstring(gradient_xml)
        fill.append(grad_fill)
    elif kwargs.get('color'):
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = kwargs.get('color')

    rectangle.line.color.rgb = RGBColor(255, 255, 255)  # Màu viền trắng
    if kwargs.get('texts'):
        text_frame = rectangle.text_frame
        text_frame.clear()  
        p = text_frame.paragraphs[0]
        p.alignment = kwargs.get('text_alignment', 1)
        run = p.add_run()
        run.text = kwargs.get('texts')
        run.font.size = kwargs.get('fontsize', Pt(6)) 
        run.font.color.rgb = kwargs.get('text_color', (0, 0, 0))
        if kwargs.get('fontweight') == "bold":
            run.font.bold = True

    if not kwargs.get('shadow'):  
        rectangle.shadow.inherit = False
        rectangle.shadow.visible = True

def create_section_rectangle(shapes, num_shapes, attribute, colors, text_colors, width_scale, **kwargs):
    left, top, width, height = attribute
    spacing = kwargs.get('spacing', Inches(0.16))
    inner_left = left
    for i in range(num_shapes):
        color = colors[i] #RGBColor(220, 220, 220)
        text_color = text_colors[i] #RGBColor(198, 198, 198)
        shape_width = width * width_scale[i]
            # color = colors[1] #RGBColor(255, 243, 205)
            # text_color = text_colors[1] #RGBColor(156, 91, 205)
            # shape_width = width*1.1    
        create_rounded_rectangle(
            shapes, (inner_left, top, shape_width, height),
            adjustment=0.5,
            color=color,
            texts=kwargs.get('texts')[i],
            fontsize=Pt(6),
            text_color=text_color,
            text_alignment=2,
            shadow=False
        )
        inner_left += spacing + shape_width
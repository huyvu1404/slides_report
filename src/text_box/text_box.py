from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_UNDERLINE

def create_text_box(shapes, texts, attribute, **kwargs):
    left, top, width, height = attribute
    textbox = shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()  # Xoá đoạn mặc định
    text_frame.word_wrap = True  # Bật xuống dòng tự động
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = texts
    run.font.size = kwargs.get('fontsize', Pt(8))
    
    if kwargs.get('fontweight') == "bold":
        run.font.bold = True
    if kwargs.get('fontstyle') == "italic":
        run.font.italic = True

    run.font.color.rgb = kwargs.get('text_color', RGBColor(0, 0, 0))
    p.alignment = kwargs.get('text_alignment', PP_ALIGN.LEFT)  # Căn trái mặc định


def create_text_box_with_url(shapes, texts, attribute, url_label=None, url=None, label_color=RGBColor(254, 219, 106),**kwargs):
    left, top, width, height = attribute
    textbox = shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    
    # Run chính
    run_main = p.add_run()
    run_main.text = texts
    run_main.font.size = kwargs.get('fontsize', Pt(12))
    run_main.font.color.rgb = kwargs.get('text_color', RGBColor(0, 0, 0))
    if kwargs.get('fontweight') == "bold":
        run_main.font.bold = True
    if kwargs.get('fontstyle') == "italic":
        run_main.font.italic = True
    if url and url_label:
        run_link = p.add_run()
        run_link.text = f" {url_label}"  
        run_link.hyperlink.address = url
        run_link.font.size = kwargs.get('fontsize', Pt(12))
        run_link.font.color.rgb = label_color
        run_link.font.underline = MSO_UNDERLINE.DASH_LINE
        if kwargs.get('fontweight') == "bold":
            run_link.font.bold = True

    p.alignment = kwargs.get('text_alignment', PP_ALIGN.LEFT)


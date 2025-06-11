from shapes import *
from text_box import *
from presentations import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from charts import create_horizontal_bar_chart
       
def create_third_slide(prs):
    _, shapes = add_slide(prs)

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    width_scale = [1, 1.1, 1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.90), Inches(0.34)
    texts = "NỘI DUNG THẢO LUẬN NỔI BẬT VỀ SHB VÀ ĐỐI THỦ"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0.14), Inches(0.89), Inches(9.82), Inches(1.2)
    texts = """Hoạt động fanpage tạo thảo luận tích cực nổi bật tại các ngân hàng Techcombank và VPBank nhờ hoạt động tương tác xúc tiến chiến dịch như minigame và livestream. Cùng với đó, choỗi hoạt động săn vé sự kiện cũng ghi nhận bàn luận sôi nổi.
Thảo luận tiêu cực của tuần vẫn xoay quanh lỗi hệ thống săn vé sự kiện trên ứng dụng và website chiến dịch, tiêu biểu đến từ VPBank và Techcombank."""

    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.1, color=RGBColor(255, 243, 205), texts=texts, fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True)
    import pandas as pd

    data = {}
    shb = pd.read_excel('./data/hbar-shb.xlsx').fillna(0)
    shb.iloc[:, 1:] = shb.iloc[:, 1:]
    data["SHB"] =  shb.iloc[:, 1:-1].values

    vcb = pd.read_excel('./data/hbar-vcb.xlsx').fillna(0)
    vcb.iloc[:, 1:] = vcb.iloc[:, 1:]
    data["Vietcombank"] = vcb.iloc[:, 1:-1].values

    tcb = pd.read_excel('./data/hbar-tcb.xlsx').fillna(0)
    tcb.iloc[:, 1:] = tcb.iloc[:, 1:]
    data["Techcombank"] = tcb.iloc[:, 1:-1].values

    mb = pd.read_excel('./data/hbar-mb.xlsx').fillna(0)
    mb.iloc[:, 1:] = mb.iloc[:, 1:]
    data["MB Bank"] = mb.iloc[:, 1:-1].values

    vp = pd.read_excel('./data/hbar-vp.xlsx').fillna(0)
    vp.iloc[:, 1:] = vp.iloc[:, 1:]
    data["VPBank"] = vp.iloc[:, 1:-1].values

    msb = pd.read_excel('./data/hbar-msb.xlsx').fillna(0)
    msb.iloc[:, 1:] = msb.iloc[:, 1:]
    data["MSB"] = msb.iloc[:, 1:-1].values

    scb = pd.read_excel('./data/hbar-scb.xlsx').fillna(0)
    scb.iloc[:, 1:] = scb.iloc[:, 1:]
    data["Sacombank"] = scb.iloc[:, 1:-1].values

    acb = pd.read_excel('./data/hbar-acb.xlsx').fillna(0)
    acb.iloc[:, 1:] = acb.iloc[:, 1:]
    data["ACB Bank"] = acb.iloc[:, 1:-1].values

    topics, width_ratios = create_horizontal_bar_chart(data)
    width_ratios[0] = 1
    print(width_ratios)
    shapes.add_picture('./img/horizontal_bar.png', Inches(0.22), Inches(2.21), Inches(9.82), Inches(3.42))
    # shapes.add_picture('./img/logo-slide3.png', Inches(1.6), Inches(2.1), width=Inches(7.91), height=Inches(0.27))
    logo_sizes = {
        "SHB": (Inches(0.43), Inches(0.12)),
        "Techcombank": (Inches(0.54), Inches(0.20)),
        "Vietcombank": (Inches(0.64), Inches(0.16)),
        "Sacombank": (Inches(0.64), Inches(0.14)),
        "MBBank": (Inches(0.33), Inches(0.16)),
        "VPBank": (Inches(0.55), Inches(0.10)),
        "ACBBank": (Inches(0.30), Inches(0.12)),
        "MSB": (Inches(0.54), Inches(0.12)),
    }
    top = Inches(2.18)
    current_left = Inches(1.88)
    for i, topic in enumerate(topics):
        topic = topic.replace(" ", "")
        
        width, height = logo_sizes.get(topic)
        left = current_left - width / 2
       
        shapes.add_picture(f'./icons/{topic}.png', left, top, width, height)
        if i < len(width_ratios) - 1:
            current_left += Inches((width_ratios[i] + width_ratios[i+1]) * 1.05 / 2)

    left, top, width, height = Inches(0.16), Inches(2.39), Inches(1.18), Inches(0.2)
    distance = Inches(0.08)
    conversations = ["BAN LÃNH ĐẠO", "CHI NHÁNH LIÊN DOANH", "CHIẾN DỊCH & SỰ KIỆN", "CHỨNG KHOÁN", "DỊCH VỤ KHÁCH HÀNG", "HOẠT ĐỘNG FANPAGE", 
          "NHÂN SỰ & TUYỂN DỤNG", "SẢN PHẨM DỊCH VỤ CÁ NHÂN", "SẢN PHẨM DỊCH VỤ DOANH NGHIỆP", "SẢN PHẨM THẺ", "THƯƠNG HIỆU CHUNG"]

    for i in range(11):
        create_rectangle(shapes, (left, top + i * (height + distance), width, height),texts=conversations[i], text_color=RGBColor(0,0,0), fontsize=Pt(5.5), color=RGBColor(254, 219, 106), shadow=False, fontweight="bold", text_alignment=2)

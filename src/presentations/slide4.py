from shapes import *
from text_box import *
from presentations import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
       
def create_fourth_slide(prs):
    _, shapes = add_slide(prs)

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    width_scale = [1, 1.1, 1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.80), Inches(0.34)
    texts = "TỔNG HỢP CHỦ ĐỀ NỔI BẬT VỀ CÁC NGÂN HÀNG"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)


    left, top, width, height = Inches(3.08), Inches(0.95), Inches(3.63), Inches(0.31)
    title = "TOP CÁC BÀI ĐĂNG TƯƠNG TÁC CAO NHẤT"
    create_rectangle(shapes, (left, top, width, height), color=RGBColor(246, 146, 0), shadow=False, fontweight="bold", texts=title, fontsize=Pt(9), text_color=RGBColor(255, 255, 255), text_alignment=2)
    
    shapes.add_picture('./img/grid_with_dashed_lines.png', Inches(0.42), Inches(1.87), width=Inches(9.48), height=Inches(3.31))
    # shapes.add_picture('./img/logo-slide4.png', Inches(0.75), Inches(1.52), width=Inches(8.74), height=Inches(0.35))
    left, top, width, height = Inches(0.03), Inches(2.37), Inches(0.35), Inches(0.24)

    colors = [RGBColor(133, 64, 134), RGBColor(221, 217, 195), RGBColor(221, 217, 195)]
    distance = Inches(0.83)
    for i in range(3):
        create_rectangle(shapes, (left, top + i * (height + distance), width, height),
                         color=colors[i], shadow=False, texts=f"{i+1}", fontsize=Pt(8), text_color=RGBColor(255, 255, 255), text_alignment=2)

    left, top, width, height = Inches(0.42), Inches(1.86), Inches(1.18), Inches(1.11)
    for i in range(3):
        for j in range(8):
            create_text_box_with_url(shapes, "Ưu đãi thẻ visa mua vé K-STAR SPARK IN VIETNAM 2025", (left + j * width, top + i * height, width, height), url_label="(URL)", url="test.com", fontsize=Pt(8), text_color=RGBColor(0, 0, 0), text_alignment=2)

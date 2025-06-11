from shapes import *
from text_box import *
from presentations import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from charts import draw_doughnut, generate_stacked_bar_chart
import pandas as pd
       
def create_first_slide(prs):
    _, shapes = add_slide(prs)

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    width_scale = [1, 1.1, 1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.90), Inches(0.34)
    texts = "TỔNG QUAN THẢO LUẬN VỀ SHB VÀ ĐỐI THỦ"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0.07), Inches(0.9), Inches(4.61), Inches(2.52)
    texts = """Tổng thảo luận toàn ngành ghi nhận giảm 44% so với tuần trước, với tổng thảo luận là 39,701 thảo luận.
    Techcombank vươn lên dẫn đầu với 32% thị phần, nổi bật với minigame Đại sứ Inspire chặng 3 và đề cập ban lãnh đạo ông Hồ Hùng Anh - Chủ tịch Techcombank thuộc tốp 3 tỷ phú Việt.
    VPBank ghi nhận giảm mạnh thị phần so với tuần trước, do sức hút từ nội dung công bố nghệ sĩ đã thu hẹp đáng kể, tiếp nối bằng choỗi chương trình săn vé sự kiện VPBank Kstar Spark2025 được quan tâm cao như sold-out pre-sale, cách săn vé. Chiến dịch Super Sinh Lời cũng được duy trì lượng nhiệt với choỗi live Podcast “Money Gym”.
    Ngoài ra, Vietcombank và MBBank cũng thu hút nhiều quan tâm đối với sự kiện hợp tác với chuỗi bán lẻ như FPT Retail và Viettel Store triển khai mô hình đại lý thanh toán."""

    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 243, 205), texts=texts, fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True)
    shapes.add_picture('./img/doughnut_side_by_side.png',Inches(4.74), Inches(1.08), Inches(4.67), Inches(2.52))
    
    title = "THỊ PHẦN THẢO LUẬN CỦA SHB VÀ ĐỐI THỦ"
    left, top, width, height = Inches(5.83), Inches(0.97), Inches(3), Inches(0.24)
    create_rectangle(shapes, (left, top, width, height), color=RGBColor(0, 112, 192), shadow=False, fontweight="bold", texts=title, fontsize=Pt(8), text_color=RGBColor(255, 255, 255), text_alignment=2)
    left, top, width, height = Inches(0.14), Inches(3.49), Inches(9.7), Inches(2.08)
    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 255, 255), shadow=True, line_color=RGBColor(227,227,227))
    
    df = pd.read_excel('./data/stacked_chart.xlsx', sheet_name='Sheet1', nrows=24, usecols='A:J').dropna(how='all').fillna(0)
    unique_topics = generate_stacked_bar_chart(df)
    print(f"Unique topics: {unique_topics}")
    
    shapes.add_picture('./img/stacked_bar_chart.png', Inches(0.29), Inches(3.57), width=Inches(9.24), height=Inches(1.77))
    left, top, width, height = Inches(0.41), Inches(3.55), Inches(3.5), Inches(0.22)
    title = "KÊNH THẢO LUẬN VỀ SHB VÀ CÁC NGÂN HÀNG ĐỐI THỦ"
    create_rectangle(shapes, (left, top, width, height), color=RGBColor(246, 146, 0), shadow=False, fontweight="bold", texts=title, fontsize=Pt(7), text_color=RGBColor(255, 255, 255), text_alignment=2)
    logo_sizes = {
        "SHB": (Inches(0.43), Inches(0.12)),
        "Techcombank": (Inches(0.54), Inches(0.20)),
        "Vietcombank": (Inches(0.64), Inches(0.16)),
        "Sacombank": (Inches(0.64), Inches(0.14)),
        "MBBank": (Inches(0.33), Inches(0.16)),
        "VPBank": (Inches(0.55), Inches(0.10)),
        "ACBBank": (Inches(0.30), Inches(0.12)),
        "MSB": (Inches(0.54), Inches(0.12)),
    }
    spacing = Inches(1.17)
    top = Inches(5.3)
    for i, topic in enumerate(unique_topics):
        topic = topic.replace(" ", "")
        print(topic)
        width, height = logo_sizes.get(topic)
        left = Inches(0.8) + i * spacing - width / 2
       
        shapes.add_picture(f'./icons/{topic}.png', left, top, width, height)


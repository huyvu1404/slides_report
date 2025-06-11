from pptx import Presentation
from pptx.util import Inches

def init():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)
    return prs

def add_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    return slide, shapes
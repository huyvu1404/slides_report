from shapes import *
from text_box import *
from presentations import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from charts import generate_sentiment_stacked_bar_chart
import pandas as pd
       
def create_second_slide(prs):
    _, shapes = add_slide(prs)

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    width_scale = [1, 1.1, 1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.01), Inches(0.34)
    texts = "SẮC THÁI THẢO LUẬN CỦA SHB VÀ ĐỐI THỦ"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0.14), Inches(0.96), Inches(9.82), Inches(1.39)
    texts = """Tổng thảo luận tích cực của các ngân hàng tăng 13% so với tuần trước. Lượng tích cực đạt kỷ lục tại Techcombank, sau đó là VPBank nhờ sức hút của nhóm bài đăng tương fanpage như minigame Đại sứ Inspire chặng 3 và live Podcast “Money Gym”.
Tổng thảo luận tiêu cực tuần này tiếp tục tăng 112% so với tuần trước. Techcombank và VPBank cũng là các ngân hàng ghi nhận tiêu cực đáng chú ý. VPBank sở hữu lượng lớn phàn nàn về việc săn vé Pre-sale khi thông báo Sold-out được công bố. Nhiều người dung phản ánh thanh toán thất bại nhưng vẫn trừ tiền, đòi ngân hàng trả lại tiền. Link săn vé cũng gặp nhiều trục trặc từ phía người sử dụng, khiến họ không thể truy cập thao tác săn vé. Tại Techcombank cũng gặp trường hợp tương tự về việc săn vé concert “Anh Trai” từ vòng quay trúng lớn trên app Techcombank Mobile."""

    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 243, 205), texts=texts, fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True)
    
    left, top, width, height = Inches(0.14), Inches(2.52), Inches(9.74), Inches(2.86)
    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 255, 255), shadow=True, line_color=RGBColor(227,227,227))
    df = pd.read_excel('./data/sentiment.xlsx', sheet_name='Sheet1', nrows=25, usecols='A, E:I').dropna(how='all').dropna(how='all', axis=1).fillna(0)
    df.columns = ['Period', 'Topic', 'Positive', 'Neutral', 'Negative']
    topics = generate_sentiment_stacked_bar_chart(df)
    print(topics)
    shapes.add_picture('./img/sentiment_stacked_bar_chart.png', Inches(0.14), Inches(2.5), Inches(9.59), Inches(2.5))
    logo_sizes = {
        "SHB": (Inches(0.43), Inches(0.12)),
        "Techcombank": (Inches(0.54), Inches(0.20)),
        "Vietcombank": (Inches(0.64), Inches(0.16)),
        "Sacombank": (Inches(0.64), Inches(0.14)),
        "MBBank": (Inches(0.33), Inches(0.16)),
        "VPBank": (Inches(0.55), Inches(0.10)),
        "ACBBank": (Inches(0.30), Inches(0.12)),
        "MSB": (Inches(0.54), Inches(0.12)),
    }
    spacing = Inches(1.17)
    top = Inches(5.01)
    for i, topic in enumerate(topics):
        topic = topic.replace(" ", "")
        print(topic)
        width, height = logo_sizes.get(topic)
        left = Inches(0.8) + i * spacing - width / 2
       
        shapes.add_picture(f'./icons/{topic}.png', left, top, width, height)
